# observepoint-revenue/tests/test_make_document.py
import json
import pathlib
import make_document

CONTENT = {
    "title": "Privacy Scan Overview",
    "subtitle": "Acme Pharma",
    "prepared_for": "Acme Pharma",
    "sections": [
        {"heading": "What we found", "body": "37 unique tags across 40 pages."},
        {"heading": "Next steps", "bullets": ["Block unapproved cookies", "Validate analytics"]},
    ],
}


def test_onepager_builds_html_with_brand(tmp_path, monkeypatch):
    monkeypatch.setattr(make_document.brand_kit, "html_to_pdf", lambda h, p: None)
    cj = tmp_path / "c.json"
    cj.write_text(json.dumps(CONTENT))
    out = tmp_path / "onepager.pdf"
    result = make_document.build("onepager", json.loads(cj.read_text()), str(out))
    assert result["engine"] is None
    produced = pathlib.Path(result["path"])
    assert produced.suffix == ".html" and produced.exists() and produced.stat().st_size > 0
    html = result["html"]
    assert "#14151A" in html.upper()            # dark theme by default
    assert "data:image/png;base64," in html      # logo embedded
    assert "Privacy Scan Overview" in html
    assert "© " in html and "ObservePoint" in html


def test_report_respects_light_override(tmp_path, monkeypatch):
    monkeypatch.setattr(make_document.brand_kit, "html_to_pdf", lambda h, p: None)
    out = tmp_path / "r.pdf"
    result = make_document.build("report", CONTENT, str(out), theme="light")
    assert "#FFFFFF" in result["html"].upper()
    assert result["theme"] == "light"


def test_letter_builds_docx_with_logo_and_copyright(tmp_path):
    from docx import Document
    out = tmp_path / "letter.docx"
    result = make_document.build("letter", CONTENT, str(out))
    assert pathlib.Path(result["path"]).suffix == ".docx"
    doc = Document(result["path"])
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "Privacy Scan Overview" in text
    assert "ObservePoint" in text
    # light theme is the default for working docs
    assert result["theme"] == "light"
    # at least one inline image (the logo) is embedded
    assert len(doc.inline_shapes) >= 1
    assert "Block unapproved cookies" in text


def test_letter_dark_theme_uses_dark_muted_color(tmp_path):
    from docx import Document
    from docx.shared import RGBColor
    out = tmp_path / "letter_dark.docx"
    result = make_document.build("letter", CONTENT, str(out), theme="dark")
    assert result["theme"] == "dark"
    doc = Document(result["path"])
    # the subtitle paragraph ("Acme Pharma") should use the dark muted color #9AA1AD
    subtitle_runs = [r for p in doc.paragraphs for r in p.runs if r.text == "Acme Pharma"]
    assert subtitle_runs, "subtitle run not found"
    assert subtitle_runs[0].font.color.rgb == RGBColor(0x9A, 0xA1, 0xAD)


def test_deck_builds_pptx_with_title_and_section_slides(tmp_path):
    from pptx import Presentation
    out = tmp_path / "deck.pptx"
    result = make_document.build("deck", CONTENT, str(out))
    assert pathlib.Path(result["path"]).suffix == ".pptx"
    assert result["theme"] == "dark"
    prs = Presentation(result["path"])
    # title slide + one slide per section
    assert len(prs.slides) == 1 + len(CONTENT["sections"])
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = "\n".join(all_text)
    assert "Privacy Scan Overview" in joined
    assert "What we found" in joined
