# observepoint-revenue/skills/branding-guide/scripts/make_document.py
"""Render a net-new ObservePoint-branded document from a content JSON.

Kinds: onepager, report  (HTML->PDF, this file)
       letter, memo       (DOCX, added in a later task)
       deck               (PPTX, added in a later task)

Claude assembles the content (judgment); this renders it (deterministic) with the
canonical brand via brand_kit. Theme defaults come from brand_kit.default_theme_for;
pass theme="dark"/"light" to override.

CLI:  python make_document.py <kind> <content.json> <out_path> [--theme dark|light]
"""
from __future__ import annotations
import datetime
import html as _html
import json
import os
import pathlib
import sys
import tempfile

import brand_kit

HTML_KINDS = {"onepager", "report"}
DOCX_KINDS = {"letter", "memo"}
PPTX_KINDS = {"deck"}


def _esc(s) -> str:
    return _html.escape("" if s is None else str(s))


def _sections_html(sections) -> str:
    out = []
    for s in sections:
        out.append(f'<h2>{_esc(s.get("heading", ""))}</h2>')
        if s.get("body"):
            out.append(f'<p>{_esc(s["body"])}</p>')
        if s.get("bullets"):
            out.append("<ul>" + "".join(f"<li>{_esc(b)}</li>" for b in s["bullets"]) + "</ul>")
    return "\n".join(out)


def render_html(content: dict, theme: str) -> str:
    year = datetime.date.today().year
    footer = content.get("footer") or brand_kit.copyright(year)
    sub = content.get("subtitle") or content.get("prepared_for") or ""
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
@import url('{brand_kit.font()["google_fonts"]}');
{brand_kit.css_vars(theme)}
*{{box-sizing:border-box}}
body{{margin:0;background:var(--op-bg,var(--op-page));color:var(--op-text);
  font-family:var(--op-font);font-size:13px;line-height:1.55;padding:48px}}
.brandbar{{display:flex;align-items:center;gap:14px;border-bottom:3px solid var(--op-accent);padding-bottom:16px}}
.brandbar img{{height:30px}}
h1{{font-weight:800;font-size:28px;margin:22px 0 2px}}
.sub{{color:var(--op-muted,var(--op-gray));font-size:13px;margin-bottom:18px}}
h2{{font-weight:800;font-size:15px;margin:20px 0 4px}}
ul{{margin:6px 0 6px 18px}}
.footer{{margin-top:40px;border-top:1px solid var(--op-border,var(--op-hairline));
  padding-top:10px;color:var(--op-muted,var(--op-gray));font-size:11px}}
</style></head><body>
<div class="brandbar"><img src="{brand_kit.logo_data_uri(theme)}" alt="ObservePoint"></div>
<h1>{_esc(content.get("title", ""))}</h1>
<div class="sub">{_esc(sub)}</div>
{_sections_html(content.get("sections", []))}
<div class="footer">{_esc(footer)}</div>
</body></html>"""


def build_docx(kind: str, content: dict, out_path: str, theme: str) -> dict:
    from docx import Document
    from docx.shared import Inches, Pt
    year = datetime.date.today().year
    family = brand_kit.font()["family"]
    ink = brand_kit.colors()["ink"]
    muted = brand_kit.theme(theme).get("muted") or brand_kit.colors()["light"]["gray"]
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name, style.font.size, style.font.color.rgb = family, Pt(10.5), brand_kit.rgbcolor(ink)
    # Letterhead: ink logo (light bg) + yellow rule
    doc.add_picture(brand_kit.logo_path(theme), width=Inches(2.0))
    rule = doc.add_paragraph()
    rule.paragraph_format.space_before = Pt(2)
    r = rule.add_run("_" * 60)
    r.font.color.rgb = brand_kit.rgbcolor(brand_kit.brand_yellow())
    # Title + subtitle
    h = doc.add_paragraph()
    hr = h.add_run(content.get("title", ""))
    hr.bold, hr.font.size, hr.font.name = True, Pt(16), family
    if content.get("subtitle") or content.get("prepared_for"):
        s = doc.add_paragraph().add_run(content.get("subtitle") or content["prepared_for"])
        s.font.size, s.font.color.rgb, s.font.name = Pt(10), brand_kit.rgbcolor(muted), family
    # Body sections
    for sec in content.get("sections", []):
        hp = doc.add_paragraph().add_run(sec.get("heading", ""))
        hp.bold, hp.font.size, hp.font.name = True, Pt(12), family
        if sec.get("body"):
            doc.add_paragraph(sec["body"])
        for b in sec.get("bullets", []):
            doc.add_paragraph(b, style="List Bullet")
    # Footer copyright
    foot = doc.add_paragraph()
    fr = foot.add_run(content.get("footer") or brand_kit.copyright(year))
    fr.font.size, fr.font.color.rgb, fr.font.name = Pt(8), brand_kit.rgbcolor(muted), family
    out = pathlib.Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return {"path": str(out), "engine": "python-docx", "theme": theme, "html": None}


def _hex_to_rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def build_pptx(kind: str, content: dict, out_path: str, theme: str) -> dict:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    t = brand_kit.theme(theme)
    bg = _hex_to_rgb(t.get("bg", t.get("page")))
    text_col = _hex_to_rgb(t["text"])
    accent = _hex_to_rgb(t["accent"])
    family = brand_kit.font()["family"]
    prs = Presentation()                       # 10 x 7.5 in default
    blank = prs.slide_layouts[6]

    def paint_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

    def add_text(slide, text, left, top, width, height, size, bold, color):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size, run.font.bold, run.font.name = Pt(size), bold, family
        run.font.color.rgb = color
        return box

    def accent_bar(slide):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.7), Inches(2.0), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    # Title slide
    s = prs.slides.add_slide(blank); paint_bg(s)
    s.shapes.add_picture(brand_kit.logo_path(theme), Inches(0.6), Inches(0.6), height=Inches(0.5))
    add_text(s, content.get("title", ""), 0.6, 1.9, 9.0, 1.5, 40, True, text_col)
    if content.get("subtitle") or content.get("prepared_for"):
        add_text(s, content.get("subtitle") or content["prepared_for"], 0.6, 3.2, 9.0, 0.8, 18, False, accent)

    # One content slide per section
    for sec in content.get("sections", []):
        cs = prs.slides.add_slide(blank); paint_bg(cs)
        add_text(cs, sec.get("heading", ""), 0.6, 0.7, 9.0, 1.0, 28, True, text_col)
        accent_bar(cs)
        body = sec.get("body", "")
        if sec.get("bullets"):
            body = (body + "\n" if body else "") + "\n".join("• " + b for b in sec["bullets"])
        if body:
            add_text(cs, body, 0.6, 2.0, 9.0, 4.5, 16, False, text_col)

    out = pathlib.Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {"path": str(out), "engine": "python-pptx", "theme": theme, "html": None}


def build(kind: str, content: dict, out_path: str, theme: str | None = None) -> dict:
    theme = theme or brand_kit.default_theme_for(kind)
    if kind in DOCX_KINDS:
        return build_docx(kind, content, out_path, theme)
    if kind in PPTX_KINDS:
        return build_pptx(kind, content, out_path, theme)
    if kind not in HTML_KINDS:
        raise ValueError(f"unknown kind {kind!r}")
    doc_html = render_html(content, theme)
    out = pathlib.Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp_html = tempfile.mkstemp(suffix=".html")
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as fh:
            fh.write(doc_html)
        engine = brand_kit.html_to_pdf(tmp_html, str(out))
    finally:
        try:
            os.unlink(tmp_html)
        except OSError:
            pass
    if engine:
        produced = str(out)
    else:                                   # no PDF engine — write the HTML beside it
        produced = str(out.with_suffix(".html"))
        pathlib.Path(produced).write_text(doc_html, encoding="utf-8")
    return {"path": produced, "engine": engine, "theme": theme, "html": doc_html}


def _main(argv) -> int:
    theme = None
    if "--theme" in argv:
        i = argv.index("--theme")
        theme = argv[i + 1]
        argv = argv[:i] + argv[i + 2:]
    if len(argv) < 3:
        sys.stderr.write("usage: make_document.py <kind> <content.json> <out_path> [--theme dark|light]\n")
        return 2
    kind, content_path, out_path = argv[0], argv[1], argv[2]
    content = json.loads(pathlib.Path(content_path).read_text())
    result = build(kind, content, out_path, theme=theme)
    print(f"wrote {result['path']} (engine={result['engine']}, theme={result['theme']})")
    return 0


if __name__ == "__main__":
    raise SystemExit(_main(sys.argv[1:]))
